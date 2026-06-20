"""Create JCA English PPTX with figures (1 per slide, widescreen)."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from PIL import Image
import os

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
figdir = os.path.join(SCRIPT_DIR, 'figures')
outdir = os.path.join(SCRIPT_DIR, 'papers')
os.makedirs(outdir, exist_ok=True)

figures = [
    {
        'file': 'fig1_price_timeseries.png',
        'label': 'Fig. 1.',
        'caption': 'Time series of eBay completed sale prices for desflurane (red), '
                   'sevoflurane (blue), and isoflurane (green) vaporizers over three years '
                   '(March 2023 to March 2026). Vertical dashed lines indicate key EU regulatory '
                   'milestones. Curved lines represent LOWESS trend estimates (fraction = 0.3). '
                   'Data source: eBay Terapeak.',
    },
    {
        'file': 'fig5_regulatory_timeline.png',
        'label': 'Fig. 2.',
        'caption': 'Anesthetic vaporizer prices mapped against the EU regulatory timeline. '
                   'Shaded regions indicate regulatory phases. Data source: eBay Terapeak.',
    },
    {
        'file': 'fig6_quarterly_trends.png',
        'label': 'Fig. 3.',
        'caption': 'Quarterly median price trends (upper panel) and sales volume (lower panel). '
                   'Data source: eBay Terapeak.',
    },
    {
        'file': 'fig2_boxplot_comparison.png',
        'label': 'Fig. 4.',
        'caption': 'Box plot comparison of vaporizer prices before and after the EU desflurane '
                   'ban (1 January 2026). Individual data points are shown as jittered dots. '
                   'Data source: eBay Terapeak.',
    },
    {
        'file': 'fig3_monthly_median.png',
        'label': 'Fig. 5.',
        'caption': 'Monthly median prices of anesthetic vaporizers on eBay. Annotations '
                   'indicate the number of transactions per month (n). Data source: eBay Terapeak.',
    },
    {
        'file': 'fig4_histograms.png',
        'label': 'Fig. 6.',
        'caption': 'Price distribution histograms for each vaporizer type, comparing pre-ban '
                   '(solid fill) and post-ban (hatched) periods. Data source: eBay Terapeak.',
    },
]


def create_pptx():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    left = Inches(1)
    top = Inches(2)
    width = slide_width - Inches(2)
    height = Inches(2)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = ('Targeted Environmental Regulation Without Observable Collateral '
              'Market Damage:\nThe EU Desflurane Ban and Secondary Market '
              'Vaporizer Prices')
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1A, 0x47, 0x6F)
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = 'Figures'
    p2.font.size = Pt(20)
    p2.font.color.rgb = RGBColor(0x55, 0x55, 0x55)
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(20)

    # Figure slides
    for fig in figures:
        fig_path = os.path.join(figdir, fig['file'])
        if not os.path.exists(fig_path):
            print(f"WARNING: {fig_path} not found, skipping")
            continue

        slide = prs.slides.add_slide(prs.slide_layouts[6])

        img_max_width = slide_width - Inches(1.5)
        img_max_height = slide_height - Inches(2.5)

        with Image.open(fig_path) as img:
            img_w, img_h = img.size
            aspect = img_w / img_h

        if aspect > (img_max_width / img_max_height):
            w = img_max_width
            h = int(w / aspect)
        else:
            h = img_max_height
            w = int(h * aspect)

        img_left = int((slide_width - w) / 2)
        img_top = Inches(0.3)

        slide.shapes.add_picture(fig_path, img_left, img_top, w, h)

        caption_top = img_top + h + Inches(0.2)
        caption_left = Inches(0.75)
        caption_width = slide_width - Inches(1.5)
        caption_height = Inches(1.5)

        txBox = slide.shapes.add_textbox(caption_left, caption_top,
                                          caption_width, caption_height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]

        run_label = p.add_run()
        run_label.text = fig['label'] + ' '
        run_label.font.size = Pt(12)
        run_label.font.bold = True

        run_caption = p.add_run()
        run_caption.text = fig['caption']
        run_caption.font.size = Pt(12)
        run_caption.font.italic = True

        p.alignment = PP_ALIGN.LEFT

    out_path = os.path.join(outdir, 'jca_figures.pptx')
    prs.save(out_path)
    print(f"JCA figures PPTX saved: {out_path}")
    return out_path


if __name__ == '__main__':
    create_pptx()
