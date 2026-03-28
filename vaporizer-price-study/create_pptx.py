#!/usr/bin/env python3
"""Create English and Japanese PPTX files with figures (1 per slide)."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

figdir = '/home/ubuntu/vaporizer_research/figures/'
outdir = '/home/ubuntu/vaporizer_research/papers/'

# Figure files and captions
figures = [
    {
        'file': 'fig1_price_timeseries.png',
        'en_label': 'Figure 1.',
        'en_caption': 'Time-series of eBay sold prices (desflurane=red, sevoflurane=blue, isoflurane=green). '
                      'Vertical dashed lines indicate EU regulatory milestones. LOWESS trend lines shown. '
                      'Note the progressive decline in the desflurane LOWESS curve. Data source: eBay Terapeak.',
        'jp_label': '\u56f31.',
        'jp_caption': 'eBay\u843d\u672d\u4fa1\u683c\u306e\u6642\u7cfb\u5217\u63a8\u79fb\uff08\u30c7\u30b9\u30d5\u30eb\u30e9\u30f3=\u8d64\u3001\u30bb\u30dc\u30d5\u30eb\u30e9\u30f3=\u9752\u3001\u30a4\u30bd\u30d5\u30eb\u30e9\u30f3=\u7dd1\uff09\u3002'
                      '\u7e26\u7834\u7dda\u306fEU\u898f\u5236\u30de\u30a4\u30eb\u30b9\u30c8\u30fc\u30f3\u3002LOWESS\u30c8\u30ec\u30f3\u30c9\u30e9\u30a4\u30f3\u4ed8\u304d\u3002'
                      '\u30c7\u30b9\u30d5\u30eb\u30e9\u30f3\u306eLOWESS\u66f2\u7dda\u306e\u6bb5\u968e\u7684\u306a\u4e0b\u964d\u306b\u6ce8\u76ee\u3002\u30c7\u30fc\u30bf\u30bd\u30fc\u30b9: eBay Terapeak\u3002',
    },
    {
        'file': 'fig2_boxplot_comparison.png',
        'en_label': 'Figure 2.',
        'en_caption': 'Pre- vs post-regulation price comparison box plots. '
                      'Note the compressed post-regulation distribution for desflurane toward lower prices. '
                      'Data source: eBay Terapeak.',
        'jp_label': '\u56f32.',
        'jp_caption': '\u898f\u5236\u524d\u5f8c\u306e\u4fa1\u683c\u6bd4\u8f03\u7bb1\u3072\u3052\u56f3\u3002\u30c7\u30b9\u30d5\u30eb\u30e9\u30f3\u306e\u898f\u5236\u5f8c\u5206\u5e03\u304c\u4f4e\u4fa1\u683c\u5074\u306b\u5727\u7e2e\u3055\u308c\u3066\u3044\u308b\u3002'
                      '\u30c7\u30fc\u30bf\u30bd\u30fc\u30b9: eBay Terapeak\u3002',
    },
    {
        'file': 'fig3_monthly_median.png',
        'en_label': 'Figure 3.',
        'en_caption': 'Monthly median price trends. A sustained decline in desflurane prices is visible from mid-2024. '
                      'Data source: eBay Terapeak.',
        'jp_label': '\u56f33.',
        'jp_caption': '\u6708\u5225\u4e2d\u592e\u5024\u4fa1\u683c\u63a8\u79fb\u30022024\u5e74\u4e2d\u9803\u304b\u3089\u30c7\u30b9\u30d5\u30eb\u30e9\u30f3\u306e\u6301\u7d9a\u7684\u4e0b\u843d\u304c\u78ba\u8a8d\u3067\u304d\u308b\u3002'
                      '\u30c7\u30fc\u30bf\u30bd\u30fc\u30b9: eBay Terapeak\u3002',
    },
    {
        'file': 'fig4_histograms.png',
        'en_label': 'Figure 4.',
        'en_caption': 'Price distribution histograms. The post-regulation desflurane distribution shifts leftward (toward lower prices). '
                      'Data source: eBay Terapeak.',
        'jp_label': '\u56f34.',
        'jp_caption': '\u4fa1\u683c\u5206\u5e03\u30d2\u30b9\u30c8\u30b0\u30e9\u30e0\u3002\u30c7\u30b9\u30d5\u30eb\u30e9\u30f3\u306e\u898f\u5236\u5f8c\u5206\u5e03\u304c\u5de6\u65b9\u5411\uff08\u4f4e\u4fa1\u683c\u5074\uff09\u306b\u30b7\u30d5\u30c8\u3057\u3066\u3044\u308b\u3002'
                      '\u30c7\u30fc\u30bf\u30bd\u30fc\u30b9: eBay Terapeak\u3002',
    },
    {
        'file': 'fig5_regulatory_timeline.png',
        'en_label': 'Figure 5.',
        'en_caption': 'EU regulatory timeline and price trends. Regulatory phases are shaded. '
                      'The stepwise price decline for desflurane is visually apparent. Data source: eBay Terapeak.',
        'jp_label': '\u56f35.',
        'jp_caption': 'EU\u898f\u5236\u30bf\u30a4\u30e0\u30e9\u30a4\u30f3\u3068\u4fa1\u683c\u63a8\u79fb\u3002\u898f\u5236\u6bb5\u968e\u3092\u7f72\u3051\u3002'
                      '\u30c7\u30b9\u30d5\u30eb\u30e9\u30f3\u306e\u6bb5\u968e\u7684\u306a\u4fa1\u683c\u4e0b\u843d\u304c\u8996\u899a\u7684\u306b\u78ba\u8a8d\u3067\u304d\u308b\u3002\u30c7\u30fc\u30bf\u30bd\u30fc\u30b9: eBay Terapeak\u3002',
    },
    {
        'file': 'fig6_quarterly_trends.png',
        'en_label': 'Figure 6.',
        'en_caption': 'Quarterly median price (top) and transaction volume (bottom) trends. '
                      'Desflurane quarterly median declined from ~$250 to ~$100. Data source: eBay Terapeak.',
        'jp_label': '\u56f36.',
        'jp_caption': '\u56db\u534a\u671f\u5225\u4e2d\u592e\u5024\u4fa1\u683c\uff08\u4e0a\uff09\u3068\u53d6\u5f15\u91cf\uff08\u4e0b\uff09\u306e\u63a8\u79fb\u3002'
                      '\u30c7\u30b9\u30d5\u30eb\u30e9\u30f3\u306e\u56db\u534a\u671f\u4e2d\u592e\u5024\u304c~$250\u304b\u3089~$100\u3078\u4e0b\u843d\u3002\u30c7\u30fc\u30bf\u30bd\u30fc\u30b9: eBay Terapeak\u3002',
    },
]


def create_pptx(lang='en'):
    """Create a PPTX with 1 figure per slide."""
    prs = Presentation()
    # Set slide dimensions to widescreen 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # Title slide
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    if lang == 'en':
        title_text = 'EU Desflurane Regulation and Secondary Market Vaporizer Prices:\nA Cross-Sectional Time-Series Analysis of eBay Completed Sales'
        subtitle_text = 'Figures and Tables'
    else:
        title_text = 'EU\u30c7\u30b9\u30d5\u30eb\u30e9\u30f3\u898f\u5236\u3068\u4e2d\u53e4\u5e02\u5834\u6c17\u5316\u5668\u4fa1\u683c\uff1a\neBay\u843d\u672d\u30c7\u30fc\u30bf\u306e\u6a2a\u65ad\u7684\u6642\u7cfb\u5217\u5206\u6790'
        subtitle_text = '\u56f3\u8868'

    # Add title text box
    left = Inches(1)
    top = Inches(2)
    width = slide_width - Inches(2)
    height = Inches(2)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1A, 0x47, 0x6F)
    p.alignment = PP_ALIGN.CENTER

    # Add subtitle
    p2 = tf.add_paragraph()
    p2.text = subtitle_text
    p2.font.size = Pt(20)
    p2.font.color.rgb = RGBColor(0x55, 0x55, 0x55)
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(20)

    # Figure slides
    for fig in figures:
        fig_path = figdir + fig['file']
        if not os.path.exists(fig_path):
            print(f"WARNING: {fig_path} not found, skipping")
            continue

        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

        label = fig['en_label'] if lang == 'en' else fig['jp_label']
        caption = fig['en_caption'] if lang == 'en' else fig['jp_caption']

        # Add figure image - centered, leaving space for caption
        # Calculate image size to fit nicely
        img_max_width = slide_width - Inches(1.5)
        img_max_height = slide_height - Inches(2.5)  # Leave space for caption

        # Get image dimensions
        from PIL import Image
        with Image.open(fig_path) as img:
            img_w, img_h = img.size
            aspect = img_w / img_h

        # Calculate scaled dimensions
        if aspect > (img_max_width / img_max_height):
            # Width-constrained
            width = img_max_width
            height = int(width / aspect)
        else:
            # Height-constrained
            height = img_max_height
            width = int(height * aspect)

        left = int((slide_width - width) / 2)
        top = Inches(0.3)

        slide.shapes.add_picture(fig_path, left, top, width, height)

        # Add caption below figure
        caption_top = top + height + Inches(0.2)
        caption_left = Inches(0.75)
        caption_width = slide_width - Inches(1.5)
        caption_height = Inches(1.5)

        txBox = slide.shapes.add_textbox(caption_left, caption_top, caption_width, caption_height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]

        # Bold label
        run_label = p.add_run()
        run_label.text = label + ' '
        run_label.font.size = Pt(12)
        run_label.font.bold = True

        # Italic caption
        run_caption = p.add_run()
        run_caption.text = caption
        run_caption.font.size = Pt(12)
        run_caption.font.italic = True

        p.alignment = PP_ALIGN.LEFT

    # Save
    if lang == 'en':
        out_path = outdir + 'vaporizer_figures_english.pptx'
    else:
        out_path = outdir + 'vaporizer_figures_japanese.pptx'

    prs.save(out_path)
    print(f"{'English' if lang == 'en' else 'Japanese'} PPTX saved: {out_path}")


if __name__ == '__main__':
    os.makedirs(outdir, exist_ok=True)
    create_pptx('en')
    create_pptx('jp')
    print("\nBoth PPTX files created successfully!")
